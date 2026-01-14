from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_ppt():
    prs = Presentation()

    # Define the content for the 7 slides
    slides_data = [
        {
            "title": "NeuroFlow",
            "subtitle": "Automated Clinical Intake & SOAP Note Generation\n\nTeam: [Your Team Name]\nBuilt for McHacks 2026",
            "layout": 0, # Title Slide
            "notes": "Hi, we are [Team Name]. We built NeuroFlow to solve the biggest bottleneck in modern healthcare: the administrative burden."
        },
        {
            "title": "The Problem: Physician Burnout",
            "content": [
                "The Bottleneck: Doctors spend 2 hours on data entry for every 1 hour of patient care.",
                "The Consequence: Burnout, shorter visits, and critical details lost.",
                "The Gap: Existing tools are just 'dictation' or require the doctor to be present."
            ],
            "layout": 1, # Title and Content
            "notes": "We found that doctors are drowning in data entry. They are becoming data clerks instead of healers. We wanted to reclaim that time."
        },
        {
            "title": "The Solution: NeuroFlow",
            "content": [
                "What it is: An intelligent agent that interviews patients BEFORE they see the doctor.",
                "Dynamic Interview: Asks smart follow-up questions based on specific symptoms.",
                "Smart Termination: Detects when it has enough info to stop.",
                "Auto-SOAP: Generates structured clinical notes instantly."
            ],
            "layout": 1,
            "notes": "NeuroFlow sits in the waiting room—digitally. It talks to the patient, gathers the history, and hands the doctor a perfect clinical note before they even walk in."
        },
        {
            "title": "Architecture & Tech Stack",
            "content": [
                "Frontend: Streamlit (Real-time Async Chat)",
                "Backend: Python 3.12 (FastAPI / AsyncIO)",
                "AI Orchestration: Backboard SDK (LLM Management)",
                "State: Custom MemoryManager & NoteTaker pipelines"
            ],
            "layout": 1,
            "notes": "Under the hood, we built a fully asynchronous event loop using Python 3.12. We integrated the Backboard SDK to handle our LLM orchestration."
        },
        {
            "title": "Key Technical Features",
            "content": [
                "🧠 Context-Aware Memory: Remembers answers to ask smart follow-ups.",
                "🛑 Semantic Termination Logic: Analyzes information density to decide when to quit.",
                "📋 Structured Output: JSON-based SOAP note generation."
            ],
            "layout": 1,
            "notes": "The hardest part was the 'Termination Logic'. We engineered the backend to analyze the chat and decide when to stop, so the patient isn't trapped in a loop."
        },
        {
            "title": "Live Demo",
            "content": [
                "1. Patient Interaction (Streamlit)",
                "2. Dynamic Follow-up Questions",
                "3. Smart Termination",
                "4. The Result: Instant SOAP Note Generation"
            ],
            "layout": 1,
            "notes": "Let's show you how it works. (Walk through the demo script. Emphasize the speed of the note generation at the end.)"
        },
        {
            "title": "Future Roadmap",
            "content": [
                "• EHR Integration (Push to Epic/Cerner)",
                "• Voice Interface (Whisper AI) for accessibility",
                "• Multi-Language Support for triage"
            ],
            "layout": 1,
            "notes": "We built the core logic this weekend. Next, we want to give NeuroFlow a voice—literally—so it can help anyone, anywhere."
        }
    ]

    for slide_info in slides_data:
        # Create slide based on layout
        slide_layout = prs.slide_layouts[slide_info["layout"]]
        slide = prs.slides.add_slide(slide_layout)

        # Set Title
        title = slide.shapes.title
        title.text = slide_info["title"]

        # Set Body / Subtitle
        if slide_info["layout"] == 0: # Title Slide
            subtitle = slide.placeholders[1]
            subtitle.text = slide_info["subtitle"]
        else: # Content Slide
            tf = slide.placeholders[1].text_frame
            tf.text = slide_info["content"][0] # First bullet
            
            for point in slide_info["content"][1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

        # Add Speaker Notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_info["notes"]

    # Save
    prs.save('NeuroFlow_Pitch.pptx')
    print("✅ PowerPoint generated: NeuroFlow_Pitch.pptx")

if __name__ == "__main__":
    create_ppt()
